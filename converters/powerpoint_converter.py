# 기존 코드를 모두 지우고 아래 코드로 교체해주세요

import os
import uuid
import tempfile
import time
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
import win32com.client
import pythoncom
import shutil

class PowerPointConverter:
    def __init__(self):
        self.temp_dir = os.path.join(tempfile.gettempdir(), 'pdf_converter')
        if not os.path.exists(self.temp_dir):
            os.makedirs(self.temp_dir)
        
        # PowerPoint 인스턴스 초기화
        self.ppt = None
        self.com_initialized = False

    def init_powerpoint(self):
        """PowerPoint 객체 초기화"""
        if self.ppt is None:
            try:
                # COM 초기화
                pythoncom.CoInitialize()
                self.com_initialized = True
                
                self.ppt = win32com.client.Dispatch("PowerPoint.Application")
                # PowerPoint 창을 표시하지 않음 (보안 정책에 따라 실패할 수 있음)
                try:
                    self.ppt.Visible = False
                except:
                    pass
                self.ppt.DisplayAlerts = False
            except Exception as e:
                if self.com_initialized:
                    pythoncom.CoUninitialize()
                    self.com_initialized = False
                raise e

    def close_powerpoint(self):
        """PowerPoint 객체 종료"""
        try:
            if self.ppt:
                self.ppt.Quit()
                self.ppt = None
        except Exception as e:
            print(f"Error closing PowerPoint: {str(e)}")
        finally:
            if self.com_initialized:
                pythoncom.CoUninitialize()
                self.com_initialized = False

    def pdf_to_ppt(self, pdf_path):
        """PDF를 PowerPoint로 변환"""
        temp_img_paths = []
        try:
            # 임시 파일 경로 생성
            output_path = os.path.join(self.temp_dir, f"{uuid.uuid4()}_output.pptx")
            
            # PDF를 이미지로 변환 (Poppler 경로 설정)
            poppler_paths = [
                r'C:\poppler-24.08.0\Library\bin',  # 기본 설치 경로
                r'C:\Program Files\poppler-24.08.0\Library\bin',  # 대체 설치 경로
                r'C:\poppler\bin'  # 수동 설치 경로
            ]

            poppler_path = None
            for path in poppler_paths:
                if os.path.exists(path):
                    poppler_path = path
                    break

            if poppler_path is None:
                # 환경 변수에서 Poppler 경로 찾기
                for path in os.environ['PATH'].split(os.pathsep):
                    if 'poppler' in path.lower() and os.path.exists(path):
                        poppler_path = path
                        break

            if poppler_path is None:
                raise Exception("Poppler not found. Please install Poppler and add it to PATH")

            images = convert_from_path(
                pdf_path,
                poppler_path=poppler_path,
                thread_count=4,
                fmt='png'
            )
            
            if not images:
                raise Exception("Failed to convert PDF to images")
            
            # 새 프레젠테이션 생성
            prs = Presentation()
            
            # 16:9 비율로 슬라이드 크기 설정
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            
            # 각 페이지를 슬라이드로 추가
            for img in images:
                # 임시 이미지 파일로 저장
                img_path = os.path.join(self.temp_dir, f"{uuid.uuid4()}.png")
                img.save(img_path, 'PNG')
                temp_img_paths.append(img_path)
                
                # 빈 슬라이드 추가
                blank_slide_layout = prs.slide_layouts[6]  # 6은 빈 슬라이드
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # 이미지 크기 조정 (슬라이드에 맞게)
                img_width, img_height = img.size
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # 이미지 비율 유지하면서 슬라이드에 맞게 크기 조정
                width_ratio = slide_width / img_width
                height_ratio = slide_height / img_height
                ratio = min(width_ratio, height_ratio)
                
                new_width = img_width * ratio
                new_height = img_height * ratio
                
                # 이미지 위치 조정 (가운데 정렬)
                left = (slide_width - new_width) / 2
                top = (slide_height - new_height) / 2
                
                # 슬라이드에 이미지 추가
                slide.shapes.add_picture(img_path, left, top, 
                                       width=new_width, 
                                       height=new_height)
            
            # 프레젠테이션 저장
            prs.save(output_path)
            
            return output_path
            
        except Exception as e:
            raise Exception(f"PDF to PowerPoint conversion failed: {str(e)}")
        finally:
            # 임시 이미지 파일 정리
            for img_path in temp_img_paths:
                try:
                    if os.path.exists(img_path):
                        os.unlink(img_path)
                except:
                    pass

    def ppt_to_pdf(self, ppt_path):
        """PowerPoint를 PDF로 변환"""
        try:
            self.init_powerpoint()
            
            try:
                # PDF 출력 경로 생성
                output_path = os.path.join(self.temp_dir, f"{uuid.uuid4()}_output.pdf")
                
                # PowerPoint 파일 열기 (ReadOnly로 열어서 창이 보이지 않도록 함)
                presentation = self.ppt.Presentations.Open(
                    FileName=os.path.abspath(ppt_path),
                    WithWindow=False,  # 창을 표시하지 않음
                    ReadOnly=True,     # 읽기 전용으로 열기
                    Untitled=True       # 제목 없는 창으로 열기
                )
                
                # PDF로 저장
                presentation.SaveAs(
                    FileName=os.path.abspath(output_path),
                    FileFormat=32  # 32 = PDF 형식
                )
                
                # 저장이 완료될 때까지 대기
                time.sleep(1)
                
                return output_path
                
            except Exception as e:
                raise Exception(f"Failed to convert PowerPoint to PDF: {str(e)}")
            finally:
                try:
                    presentation.Close()
                except:
                    pass
        except Exception as e:
            raise Exception(f"PowerPoint initialization failed: {str(e)}")
        finally:
            self.close_powerpoint()

    def cleanup(self):
        """임시 파일 정리"""
        try:
            if os.path.exists(self.temp_dir):
                shutil.rmtree(self.temp_dir)
        except Exception as e:
            print(f"Error cleaning up temporary files: {str(e)}")